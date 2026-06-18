"""
PPTX extractor for DOCINTEL.
Extracts text and metadata from PPTX files using python-pptx.
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from docintel.common.exceptions import ExtractionError
from docintel.common.logging import get_logger
from docintel.common.constants import SOURCE_LAYER_NATIVE

logger = get_logger(__name__)


class PPTXExtractor:
    """
    Extracts text and metadata from PPTX files.
    Uses python-pptx for reading PowerPoint presentations.
    """
    
    def __init__(self):
        """Initialize the PPTX extractor."""
        pass
    
    def extract(self, file_path: str) -> dict:
        """
        Extract text and metadata from a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Dictionary with 'text', 'metadata', and 'source_layer'
        """
        try:
            prs = Presentation(file_path)
            
            # Extract metadata
            metadata = self._extract_metadata(prs)
            
            # Extract text from all slides
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = self._extract_slide_text(slide, slide_num)
                if slide_text:
                    text_parts.append(f"\n[Slide {slide_num}]\n{slide_text}")
            
            full_text = "\n".join(text_parts)
            
            logger.info(
                "PPTX extraction complete",
                extra={
                    "file": file_path,
                    "slides": len(prs.slides),
                    "total_chars": len(full_text),
                },
            )
            
            return {
                "text": full_text,
                "metadata": metadata,
                "source_layer": SOURCE_LAYER_NATIVE,
            }
            
        except Exception as e:
            error_msg = f"Failed to extract PPTX: {str(e)}"
            logger.error(error_msg, extra={"file": file_path}, exc_info=True)
            raise ExtractionError(error_msg) from e
    
    def _extract_slide_text(self, slide, slide_num: int) -> str:
        """
        Extract text from a slide.
        
        Args:
            slide: python-pptx slide object
            slide_num: Slide number for logging
            
        Returns:
            Extracted text
        """
        try:
            text_parts = []
            
            # Extract from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
                
                # Extract from tables
                if shape.has_table:
                    table_text = self._extract_table_text(shape.table)
                    if table_text:
                        text_parts.append(f"[Table]\n{table_text}")
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"Could not extract slide {slide_num}: {str(e)}")
            return ""
    
    def _extract_table_text(self, table) -> str:
        """
        Extract text from a PPTX table.
        
        Args:
            table: python-pptx table object
            
        Returns:
            Extracted table text
        """
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            return "\n".join(rows)
        except Exception as e:
            logger.warning(f"Could not extract table: {str(e)}")
            return ""
    
    def _extract_metadata(self, prs: Presentation) -> dict:
        """
        Extract metadata from a PPTX presentation.
        
        Args:
            prs: python-pptx presentation object
            
        Returns:
            Dictionary of metadata
        """
        metadata = {}
        
        try:
            # Get core properties
            core_props = prs.core_properties
            metadata["title"] = core_props.title or ""
            metadata["author"] = core_props.author or ""
            metadata["subject"] = core_props.subject or ""
            metadata["created"] = str(core_props.created) if core_props.created else ""
            metadata["modified"] = str(core_props.modified) if core_props.modified else ""
            
            # Get slide dimensions
            metadata["slide_width"] = prs.slide_width
            metadata["slide_height"] = prs.slide_height
            
            # Count slides
            metadata["slide_count"] = len(prs.slides)
            
        except Exception as e:
            logger.warning(f"Could not extract all PPTX metadata: {str(e)}")
        
        return metadata
    
    def extract_slide(self, file_path: str, slide_number: int) -> dict:
        """
        Extract text from a specific slide.
        
        Args:
            file_path: Path to the PPTX file
            slide_number: 1-based slide number
            
        Returns:
            Dictionary with 'text' and 'metadata'
        """
        try:
            prs = Presentation(file_path)
            
            if slide_number < 1 or slide_number > len(prs.slides):
                raise ValueError(f"Slide {slide_number} out of range (1-{len(prs.slides)})")
            
            slide = prs.slides[slide_number - 1]
            text = self._extract_slide_text(slide, slide_number)
            
            metadata = {
                "slide_number": slide_number,
                "source_format": "PPTX",
            }
            
            logger.info(
                "PPTX slide extraction complete",
                extra={"file": file_path, "slide": slide_number, "chars": len(text)},
            )
            
            return {
                "text": text,
                "metadata": metadata,
                "source_layer": SOURCE_LAYER_NATIVE,
            }
            
        except Exception as e:
            error_msg = f"Failed to extract PPTX slide: {str(e)}"
            logger.error(error_msg, extra={"file": file_path}, exc_info=True)
            raise ExtractionError(error_msg) from e
    
    def get_slide_count(self, file_path: str) -> int:
        """
        Get the number of slides in a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Number of slides
        """
        try:
            prs = Presentation(file_path)
            return len(prs.slides)
        except Exception as e:
            error_msg = f"Failed to count slides: {str(e)}"
            logger.error(error_msg, extra={"file": file_path})
            raise ExtractionError(error_msg) from e


# Global PPTX extractor instance
_pptx_extractor: Optional[PPTXExtractor] = None


def get_pptx_extractor() -> PPTXExtractor:
    """Get the global PPTX extractor instance."""
    global _pptx_extractor
    if _pptx_extractor is None:
        _pptx_extractor = PPTXExtractor()
    return _pptx_extractor