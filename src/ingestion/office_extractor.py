"""Office document extractor for DOCX/XLSX/PPTX using standard libraries."""
import logging
from pathlib import Path
from typing import Optional

from src.common.exceptions import IngestionError, ParseError

logger = logging.getLogger(__name__)


class OfficeExtractor:
    """
    Office document extractor.
    Provides structured access to DOCX, XLSX, and PPTX files using standard libraries.
    """
    
    def extract(self, file_path: Path, page_count: Optional[int] = None) -> dict:
        """
        Extract text and metadata from Office document.
        
        Args:
            file_path: Path to Office document
            page_count: Optional known page count
            
        Returns:
            Dictionary with extracted text and metadata
        """
        extension = file_path.suffix.lower()
        
        try:
            if extension == ".docx":
                return self._extract_docx(file_path, page_count)
            elif extension == ".xlsx":
                return self._extract_xlsx(file_path, page_count)
            elif extension == ".pptx":
                return self._extract_pptx(file_path, page_count)
            else:
                raise IngestionError(f"Unsupported Office format: {extension}")
        
        except Exception as e:
            logger.error(f"Office extraction failed: {e}", exc_info=True)
            raise ParseError(f"Office extraction failed: {e}") from e
    
    def _extract_docx(self, file_path: Path, page_count: Optional[int] = None) -> dict:
        """
        Extract text from DOCX file using python-docx.
        
        Args:
            file_path: Path to DOCX file
            page_count: Optional known page count
            
        Returns:
            Extraction results
        """
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(file_path)
            text_parts = []
            headings = []
            tables = []
            
            # Extract paragraphs and headings
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
                    
                    # Check if it's a heading
                    if para.style and para.style.name.startswith("Heading"):
                        headings.append({
                            "level": para.style.name,
                            "text": para.text,
                        })
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            text = "\n".join(text_parts)
            
            return {
                "text": text,
                "blocks": [],
                "low_confidence_regions": [],
                "page_count": page_count or 1,
                "extraction_method": "python-docx",
                "headings": headings,
                "tables": tables,
            }
        
        except ImportError:
            raise IngestionError("python-docx not available for DOCX extraction")
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise
    
    def _extract_xlsx(self, file_path: Path, page_count: Optional[int] = None) -> dict:
        """
        Extract text from XLSX file using openpyxl.
        
        Args:
            file_path: Path to XLSX file
            page_count: Optional known page count
            
        Returns:
            Extraction results
        """
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, read_only=True)
            text_parts = []
            sheets_info = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n[Sheet: {sheet_name}]\n")
                
                sheet_rows = 0
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    text_parts.append(row_text)
                    sheet_rows += 1
                
                sheets_info.append({
                    "name": sheet_name,
                    "rows": sheet_rows,
                })
            
            wb.close()
            
            return {
                "text": "\n".join(text_parts),
                "blocks": [],
                "low_confidence_regions": [],
                "page_count": len(sheets_info),
                "extraction_method": "openpyxl",
                "sheets": sheets_info,
            }
        
        except ImportError:
            raise IngestionError("openpyxl not available for XLSX extraction")
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            raise
    
    def _extract_pptx(self, file_path: Path, page_count: Optional[int] = None) -> dict:
        """
        Extract text from PPTX file using python-pptx.
        
        Args:
            file_path: Path to PPTX file
            page_count: Optional known page count
            
        Returns:
            Extraction results
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            slides_info = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n[Slide {i}]\n")
                
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                
                text_parts.extend(slide_texts)
                
                slides_info.append({
                    "number": i,
                    "text_blocks": len(slide_texts),
                })
            
            return {
                "text": "\n".join(text_parts),
                "blocks": [],
                "low_confidence_regions": [],
                "page_count": len(slides_info),
                "extraction_method": "python-pptx",
                "slides": slides_info,
            }
        
        except ImportError:
            raise IngestionError("python-pptx not available for PPTX extraction")
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise


# Global Office extractor instance
office_extractor = OfficeExtractor()